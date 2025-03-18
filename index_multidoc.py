import os
from pathlib import Path
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, Settings
from llama_index.core.schema import Document
from llama_index.llms.ollama import Ollama
from llama_index.embeddings.ollama import OllamaEmbedding
from llama_index.readers.file import PyMuPDFReader, UnstructuredReader
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import tempfile
import shutil

def setup_models():
    """Setup LLM and embedding model for indexing"""
    llm = Ollama(model="llama3.2", temperature=0.1)
    embed_model = OllamaEmbedding(model_name="nomic-embed-text", base_url="http://localhost:11434")
    
    Settings.llm = llm
    Settings.embed_model = embed_model
    return llm, embed_model

def extract_text_from_image(image_path):
    """Extract text from an image using OCR"""
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Error extracting text from image {image_path}: {e}")
        return ""

def extract_text_from_word(docx_path):
    """Extract text from a Word document"""
    try:
        doc = docx.Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error extracting text from Word document {docx_path}: {e}")
        return ""

def extract_text_from_excel(excel_path):
    """Extract text from an Excel spreadsheet"""
    try:
        workbook = openpyxl.load_workbook(excel_path, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}")
            
            for row in sheet.iter_rows():
                row_texts = []
                for cell in row:
                    if cell.value is not None:
                        row_texts.append(str(cell.value))
                if row_texts:
                    text_parts.append(" | ".join(row_texts))
        
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from Excel file {excel_path}: {e}")
        return ""

def extract_text_from_powerpoint(pptx_path):
    """Extract text from a PowerPoint presentation"""
    try:
        presentation = Presentation(pptx_path)
        text_parts = []
        
        for i, slide in enumerate(presentation.slides):
            text_parts.append(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from PowerPoint file {pptx_path}: {e}")
        return ""

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF using PyMuPDF"""
    try:
        doc = fitz.open(pdf_path)
        text_parts = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text_parts.append(page.get_text())
        
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from PDF {pdf_path}: {e}")
        return ""

def process_documents(input_dir, output_dir="./index_storage"):
    """Process all supported documents and create an index"""
    # Setup models
    llm, embed_model = setup_models()
    
    # Create temp directory for processed files
    temp_dir = tempfile.mkdtemp()
    print(f"Created temporary directory: {temp_dir}")
    
    try:
        documents = []
        
        # Process all files in the input directory
        for root, _, files in os.walk(input_dir):
            for file in files:
                file_path = os.path.join(root, file)
                file_ext = os.path.splitext(file)[1].lower()
                print(f"Processing: {file_path}")
                
                # Process based on file extension
                if file_ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
                    # Image files - extract text with OCR
                    text = extract_text_from_image(file_path)
                    if text:
                        doc = Document(
                            text=text,
                            metadata={"source": file_path, "file_type": "image"}
                        )
                        documents.append(doc)
                
                elif file_ext == '.docx':
                    # Word documents
                    text = extract_text_from_word(file_path)
                    if text:
                        doc = Document(
                            text=text,
                            metadata={"source": file_path, "file_type": "word"}
                        )
                        documents.append(doc)
                
                elif file_ext in ['.xlsx', '.xls']:
                    # Excel spreadsheets
                    text = extract_text_from_excel(file_path)
                    if text:
                        doc = Document(
                            text=text,
                            metadata={"source": file_path, "file_type": "excel"}
                        )
                        documents.append(doc)
                
                elif file_ext in ['.pptx', '.ppt']:
                    # PowerPoint presentations
                    text = extract_text_from_powerpoint(file_path)
                    if text:
                        doc = Document(
                            text=text,
                            metadata={"source": file_path, "file_type": "powerpoint"}
                        )
                        documents.append(doc)
                
                elif file_ext == '.pdf':
                    # PDF files
                    text = extract_text_from_pdf(file_path)
                    if text:
                        doc = Document(
                            text=text,
                            metadata={"source": file_path, "file_type": "pdf"}
                        )
                        documents.append(doc)
                
                elif file_ext == '.txt':
                    # Text files
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                    doc = Document(
                        text=text,
                        metadata={"source": file_path, "file_type": "text"}
                    )
                    documents.append(doc)
        
        print(f"Processed {len(documents)} documents")
        
        # Create the index
        if documents:
            from llama_index.core.node_parser import SentenceWindowNodeParser
            
            # Create parser for chunking documents
            node_parser = SentenceWindowNodeParser.from_defaults(
                window_size=3,
                window_metadata_key="window",
                original_text_metadata_key="original_text"
            )
            
            # Create and save the index
            index = VectorStoreIndex.from_documents(
                documents,
                node_parser=node_parser
            )
            
            index.storage_context.persist(output_dir)
            print(f"Index created and saved to {output_dir}")
        else:
            print("No documents were processed. Index not created.")
    
    finally:
        # Clean up temporary directory
        shutil.rmtree(temp_dir)
        print(f"Removed temporary directory: {temp_dir}")

if __name__ == "__main__":
    # Directory containing documents to index
    input_directory = "./data"
    
    # Directory where the index will be stored
    output_directory = "./index_storage"
    
    process_documents(input_directory, output_directory)